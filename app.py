import streamlit as st
from groq import Groq
import io
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

# Helper Function: Create Word Document
def create_docx(content):
    doc = Document()
    doc.add_heading('Lesson Plan', 0)
    for paragraph in content.split('\n\n'):
        if paragraph.strip():
            doc.add_paragraph(paragraph.strip())
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

# Smart Helper Function: Parse AI Text into a 10-12 Slide Presentation
def create_smart_pptx(slide_text):
    prs = Presentation()
    
    # Split the AI output by the word "Slide" to separate individual slides
    raw_slides = slide_text.split('Slide ')
    
    for raw_slide in raw_slides:
        if not raw_slide.strip():
            continue
            
        # Add a standard Title + Bullet Content slide layout
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = raw_slide.strip().split('\n')
        
        # Line 0 usually contains the Slide number and Slide Title (e.g., "1: Welcome to the Lesson")
        slide_title = lines[0].split(':', 1)[-1].strip() if ':' in lines[0] else lines[0]
        slide.shapes.title.text = slide_title
        
        # Clean up remaining lines to act as slide body bullet points
        body_lines = [line.strip().replace('- ', '').replace('* ', '') for line in lines[1:] if line.strip()]
        body_text = "\n".join(body_lines)
        
        slide.placeholders[1].text = body_text
        
    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()

# 1. Page Configuration
st.set_page_config(
    page_title="Teacher Bot",
    page_icon="🇦🇪",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide Streamlit Header, Main Menu, and GitHub Deploy/Source Links
hide_streamlit_style = """
    <style>
    #MainMenu {visibility: hidden;}
    header {visibility: hidden;}
    footer {visibility: hidden;}
    .viewerBadge_container__1QSob {display: none !important;}
    </style>
"""
st.markdown(hide_streamlit_style, unsafe_allow_html=True)


# 2. Secure API Key Management (Groq)
if "GROQ_API_KEY" in st.secrets:
    client = Groq(api_key=st.secrets["GROQ_API_KEY"])
else:
    st.sidebar.warning("Groq API Key not found in Secrets. Please provide it below:")
    user_key = st.sidebar.text_input("Enter Groq API Key", type="password")
    if user_key:
        client = Groq(api_key=user_key)
    else:
        client = None

# Initialize Session State variables
if "lesson_plan" not in st.session_state:
    st.session_state.lesson_plan = ""
if "student_slides_text" not in st.session_state:
    st.session_state.student_slides_text = ""
if "component_output" not in st.session_state:
    st.session_state.component_output = ""
if "active_component_name" not in st.session_state:
    st.session_state.active_component_name = "Framework Component Output"

# Helper function to query Groq
def call_groq(prompt_text):
    if not client:
        return "⚠️ Error: Groq API client is not configured. Please add your key."
    try:
        chat_completion = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt_text}],
            model="llama-3.3-70b-versatile",
        )
        return chat_completion.choices[0].message.content
    except Exception as e:
        return f"Error communicating with Groq API: {str(e)}"

# 3. Sidebar UI (Lesson Parameters)
with st.sidebar:
    st.markdown("### 🇦🇪 Teacher Bot")
    st.caption("UAE-Aligned Smart Lesson Builder")
    st.divider()
    
    st.markdown("#### **Lesson Parameters**")
    grade = st.selectbox("Grade Level / Year", [f"Grade {i}" for i in range(1, 13)] + [f"Year {i}" for i in range(1, 14)])
    subject = st.selectbox("Subject / Course", ["Mathematics", "Science", "English Language Arts", "Social Studies", "Moral Education", "Islamic Studies", "Arabic Language", "Computing / ICT"])
    topic = st.text_input("Lesson Topic", placeholder="e.g., Sustainable Eco-Cities")
    objectives = st.text_area("Learning Objectives (LOs)", placeholder="What should students know or be able to do?")
    
    with st.expander("Advanced Settings"):
        time_duration = st.slider("Lesson Duration (Minutes)", 30, 120, 60, step=5)
        differentiation = st.multiselect("Differentiation Focus", ["SEN Support", "G&T (Gifted & Talented)", "ELL / Language Support"], default=["SEN Support"])

    st.write("")
    generate_btn = st.button("⚡ Generate Full Lesson Plan", type="primary", use_container_width=True)
    
    # NEW SLIDE GENERATOR BUTTON RIGHT IN THE SIDEBAR!
    generate_slides_btn = st.button("🖼️ Generate Student PPT Slides", type="secondary", use_container_width=True)

# Insert this right before: col_workspace, col_tools = st.columns([1.1, 1.3])

# Mobile Sidebar Rescue Banner
st.markdown(
    """
    <style>
    @media (max-width: 768px) {
        .mobile-hint {
            background-color: #f0f2f6;
            padding: 10px;
            border-radius: 5px;
            margin-bottom: 15px;
            text-align: center;
            font-size: 14px;
        }
    }
    @media (min-width: 769px) {
        .mobile-hint { display: none; }
    }
    </style>
    <div class="mobile-hint">
        📱 <b>Using a phone?</b> Click the small arrow <b>&gt;</b> in the top-left corner if you need to change Lesson Parameters!
    </div>
    """, 
    unsafe_allow_html=True
)

# 4. Layout: Merged into 2 clean columns
col_workspace, col_tools = st.columns([1.1, 1.3])

# --- COLUMN 1: MAIN WORKSPACE (LESSON PLAN OR PPT SLIDES VIEW) ---
with col_workspace:
    # Trigger 1: Handle Lesson Plan Generation
    if generate_btn:
        if not topic or not objectives:
            st.error("⚠️ Please fill in both the Topic and Learning Objectives to start.")
        else:
            with st.spinner("Generating full framework via Groq..."):
                base_prompt = f"""
                You are an expert curriculum designer specializing in the UAE Framework.
                Generate a comprehensive, structured teacher lesson plan based on these parameters:
                - Grade: {grade}, Subject: {subject}, Topic: {topic}, Objectives: {objectives}, Duration: {time_duration} mins.
                Organize clearly into standard sections.
                """
                st.session_state.lesson_plan = call_groq(base_prompt)
                # Clear active slides text so we view the plan
                st.session_state.student_slides_text = ""

    # Trigger 2: Handle Student Classroom Slides Generation
    if generate_slides_btn:
        if not topic or not objectives:
            st.error("⚠️ Please fill in both the Topic and Learning Objectives to start.")
        else:
            with st.spinner("Drafting 10-12 Interactive Student Slides..."):
                slides_prompt = f"""
                Create a slide-by-slide classroom presentation script tailored for students based on:
                - Grade: {grade}, Subject: {subject}, Topic: {topic}, Objectives: {objectives}
                
                You must generate between 10 to 12 distinct slides. Format your entire output STRICTLY like this blueprint example, using the word 'Slide X: Title' to split slides cleanly:
                
                Slide 1: Lesson Title & Hook
                - Welcome to our lesson on {topic}!
                - Here is our mystery question for today...
                
                Slide 2: Our Learning Goals
                - By the end of today, we will understand...
                - We will be able to...
                
                Follow this structure for all 10-12 slides covering Vocabulary, Introduction, Guided Task, Independent Task, Group Challenge, Quiz, and Reflection. Do not include extra conversational text outside this structure.
                """
                st.session_state.student_slides_text = call_groq(slides_prompt)
                # Clear active lesson plan view so we look at slides
                st.session_state.lesson_plan = ""

    # DISPLAY RESULTS DYNAMICALLY IN WORKSPACE
    if st.session_state.lesson_plan:
        st.subheader("📝 Main Workspace: Lesson Plan")
        st.markdown(st.session_state.lesson_plan)
        st.divider()
        
        btn_col1, btn_col2 = st.columns(2)
        with btn_col1:
            st.download_button("📥 Download Word (.docx)", data=create_docx(st.session_state.lesson_plan), file_name=f"Lesson_Plan_{topic.replace(' ', '_')}.docx", use_container_width=True)
        with btn_col2:
            st.download_button("📥 Download Text (.txt)", data=st.session_state.lesson_plan, file_name=f"Lesson_Plan_{topic.replace(' ', '_')}.txt", use_container_width=True)

    elif st.session_state.student_slides_text:
        st.subheader("🖼️ Main Workspace: Classroom Presentation Preview")
        st.markdown(st.session_state.student_slides_text)
        st.divider()
        
        # Provide the real 10-12 multi-slide download!
        st.download_button(
            label="📥 Download Classroom PPT (.pptx)", 
            data=create_smart_pptx(st.session_state.student_slides_text), 
            file_name=f"Classroom_Slides_{topic.replace(' ', '_')}.pptx", 
            use_container_width=True,
            type="primary"
        )
    else:
        st.subheader("📝 Main Workspace")
        st.info("👋 Fill out the sidebar and click 'Generate Full Lesson Plan' to view your teacher guide, or click 'Generate Student PPT Slides' to build classroom presentation slides!")


# --- COLUMN 2: MERGED FRAMEWORK BLOCKS & COPILOT TOOLS ---
with col_tools:
    st.subheader("🛠️ Framework Blocks")
    st.caption("Click a button to generate a specific sub-framework element below")
    
    comp_col1, comp_col2, comp_col3 = st.columns(3)
    selected_component = None
    component_instruction = ""
    
    with comp_col1:
        if st.button("📋 Standards", use_container_width=True, help="Matching Standards"):
            selected_component = "📋 Matching Standards"
            component_instruction = "Generate curriculum benchmarks and alignment metrics."
        if st.button("💡 Big Idea", use_container_width=True, help="Essential Question"):
            selected_component = "💡 Essential Question / Big Idea"
            component_instruction = "Develop inquiry-driven core statements."
        if st.button("🌍 Real-Life", use_container_width=True, help="Real-Life Connections"):
            selected_component = "🌍 Real-Life Connections"
            component_instruction = "Provide real-world application connections."
        if st.button("🔗 Interdisc.", use_container_width=True, help="Interdisciplinary"):
            selected_component = "🔗 Interdisciplinary Connections"
            component_instruction = "Link topic seamlessly with another subject."
        if st.button("🇦🇪 UAE National", use_container_width=True, help="UAE National Identity"):
            selected_component = "🇦🇪 UAE National Identity"
            component_instruction = "Incorporate UAE history, culture, or vision initiatives."

    with comp_col2:
        if st.button("🌟 Profiles", use_container_width=True, help="Learner Profiles"):
            selected_component = "🌟 Learner Profile Attributes"
            component_instruction = "Outline targeted student-centric profile descriptions."
        if st.button("📝 Pre-Assess", use_container_width=True, help="Pre-Assessment (5 MCQ)"):
            selected_component = "📝 Pre-Assessment (5 MCQs)"
            component_instruction = "Draft a 5-question diagnostic MCQ assessment with an answer key."
        if st.button("🔥 Ind. Task", use_container_width=True, help="Independent Task"):
            selected_component = "🔥 Independent Learning Task"
            component_instruction = "Design a tiered individual assignment (Mild, Medium, Spicy)."
        if st.button("👥 Group Prac.", use_container_width=True, help="Group Practice"):
            selected_component = "👥 Group Practice (Max 20 mins)"
            component_instruction = "Construct a high-collaboration 20-minute group challenge."

    with comp_col3:
        if st.button("🚀 PBL Ideas", use_container_width=True, help="Project-Based Learning"):
            selected_component = "🚀 PBL Ideas"
            component_instruction = "Propose creative long-term Project-Based Learning prompts."
        if st.button("♿ SOD Support", use_container_width=True, help="Students of Determination"):
            selected_component = "♿ Students of Determination (SOD)"
            component_instruction = "Provide specific accommodations and scaffolded tools for SOD requirements."
        if st.button("🎫 Exit Ticket", use_container_width=True, help="Exit Ticket"):
            selected_component = "🎫 Exit Ticket"
            component_instruction = "Formulate a quick reflective exit check prompt."
        if st.button("🎯 Post-Assess", use_container_width=True, help="Post-Assessment (3 MCQ)"):
            selected_component = "🎯 Post-Assessment (3 MCQs)"
            component_instruction = "Compose a 3-question summative evaluation tool with an answer key."

    if selected_component and component_instruction:
        if not topic or not objectives:
            st.error("⚠️ Make sure you have entered a Topic and Objectives in the sidebar first!")
        else:
            with st.spinner(f"Generating {selected_component}..."):
                prompt = f"Task: Generate '{selected_component}' for Grade: {grade}, Subject: {subject}, Topic: {topic}.\n{component_instruction}\nReturn only plain structured text."
                st.session_state.component_output = call_groq(prompt)
                st.session_state.active_component_name = selected_component

    st.divider()
    st.markdown(f"##### 🔲 {st.session_state.active_component_name}")
    if st.session_state.component_output:
        st.code(st.session_state.component_output, language="markdown")
    else:
        st.caption("Click any framework block above to view text snippets here.")

    # --- INTEGRATED AI COPILOT REFINEMENT PANEL ---
    st.divider()
    st.subheader("🤖 AI Copilot")
    st.caption("Apply quick refinements or custom instructions to your left Main Workspace Plan")
    
    cop_col1, cop_col2, cop_col3 = st.columns(3)
    with cop_col1:
        add_quiz = st.button("📝 Add Extra Quiz Items", use_container_width=True)
    with cop_col2:
        more_sen = st.button("🤝 Deepen Support Strategies", use_container_width=True)
    with cop_col3:
        time_breakdown = st.button("⏱️ Create Detailed Pacing", use_container_width=True)
    
    chip_instruction = ""
    if add_quiz: chip_instruction = "Expand the assessment metrics by attaching additional custom quiz tracking options."
    if more_sen: chip_instruction = "Deepen inclusion scaffolding adjustments across all target lesson blocks."
    if time_breakdown: chip_instruction = "Provide an exact, comprehensive minute-by-minute pace structure tracking how to guide this entire agenda."

    user_instruction = st.text_input("💬 Custom adjustments to Main Plan...", placeholder="e.g., Translate vocabulary to Arabic...")
    submit_instruction = st.button("Apply to Main Plan", use_container_width=True)
    
    final_instruction = user_instruction if submit_instruction else chip_instruction
    
    if final_instruction:
        if not st.session_state.lesson_plan and not st.session_state.student_slides_text:
            st.error("Please generate either a Lesson Plan or PPT Slides in the workspace column before requesting edits.")
        else:
            with st.spinner("Copilot adapting layout..."):
                active_text = st.session_state.lesson_plan if st.session_state.lesson_plan else st.session_state.student_slides_text
                refinement_prompt = f"You are editing a master curriculum draft.\nCURRENT DRAFT:\n---\n{active_text}\n---\nUSER AMENDMENT DIRECTION: \"{final_instruction}\"\nReturn the fully reconstructed integrated schema text cleanly."
                
                if st.session_state.lesson_plan:
                    st.session_state.lesson_plan = call_groq(refinement_prompt)
                else:
                    st.session_state.student_slides_text = call_groq(refinement_prompt)
                st.rerun()
